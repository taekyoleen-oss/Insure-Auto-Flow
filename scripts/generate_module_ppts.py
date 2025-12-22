#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
모듈별 PPT 생성 스크립트
하나의 PPT 파일에 모든 모듈을 포함합니다. 각 모듈은 하나의 슬라이드로 표현됩니다.
"""

import json
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from openai import OpenAI
import re
import platform

# OpenAI API 키 설정 (환경 변수에서 가져오기)
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")

if OPENAI_API_KEY:
    client = OpenAI(api_key=OPENAI_API_KEY)
else:
    client = None
    print("경고: OPENAI_API_KEY가 설정되지 않았습니다. AI 설명을 생성할 수 없습니다.")

# Shape 타입 제외 목록
EXCLUDED_TYPES = ['TextBox', 'GroupBox']

def get_download_folder():
    """사용자의 다운로드 폴더 경로 반환"""
    if platform.system() == "Windows":
        import winreg
        sub_key = r'SOFTWARE\Microsoft\Windows\CurrentVersion\Explorer\Shell Folders'
        downloads_guid = '{374DE290-123F-4565-9164-39C4925E467B}'
        with winreg.OpenKey(winreg.HKEY_CURRENT_USER, sub_key) as key:
            location = winreg.QueryValueEx(key, downloads_guid)[0]
        return location
    elif platform.system() == "Darwin":  # macOS
        return os.path.join(os.path.expanduser("~"), "Downloads")
    else:  # Linux
        return os.path.join(os.path.expanduser("~"), "Downloads")

def sanitize_filename(name):
    """파일명으로 사용할 수 없는 문자 제거"""
    name = re.sub(r'[<>:"/\\|?*]', '_', name)
    name = name.replace(' ', '_')
    return name

def get_input_data_info(module, all_modules, connections):
    """모듈의 입력 데이터 정보 추출"""
    inputs = module.get('inputs', [])
    if not inputs:
        return "입력 데이터 없음"
    
    input_info = []
    for input_port in inputs:
        # 연결된 모듈 찾기
        connected_module = None
        for conn in connections:
            if conn.get('to', {}).get('moduleId') == module.get('id') and \
               conn.get('to', {}).get('portName') == input_port.get('name'):
                from_module_id = conn.get('from', {}).get('moduleId')
                connected_module = next((m for m in all_modules if m.get('id') == from_module_id), None)
                break
        
        port_type = input_port.get('type', 'unknown')
        if connected_module:
            module_name = connected_module.get('name', connected_module.get('type'))
            output_data = connected_module.get('outputData')
            
            if output_data and output_data.get('type') == 'DataPreview':
                columns = output_data.get('columns', [])
                row_count = output_data.get('totalRowCount', 0)
                input_info.append(f"• {input_port['name']} ({port_type}): {module_name}에서 전달\n  - 형태: {row_count}행 × {len(columns)}열")
                if columns:
                    col_names = [col.get('name', '') for col in columns[:5]]
                    input_info.append(f"  - 주요 컬럼: {', '.join(col_names)}{'...' if len(columns) > 5 else ''}")
            else:
                input_info.append(f"• {input_port['name']} ({port_type}): {module_name}에서 전달")
        else:
            input_info.append(f"• {input_port['name']} ({port_type}): 연결되지 않음")
    
    return '\n'.join(input_info) if input_info else "입력 데이터 없음"

def get_output_data_info(module):
    """모듈의 출력 데이터 정보 추출"""
    output_data = module.get('outputData')
    outputs = module.get('outputs', [])
    
    if not outputs:
        return "출력 데이터 없음"
    
    output_info = []
    for output_port in outputs:
        port_type = output_port.get('type', 'unknown')
        port_name = output_port.get('name', 'unknown')
        
        if output_data:
            if output_data.get('type') == 'DataPreview':
                columns = output_data.get('columns', [])
                row_count = output_data.get('totalRowCount', 0)
                output_info.append(f"• {port_name} ({port_type}): 데이터 테이블\n  - 형태: {row_count}행 × {len(columns)}열")
                if columns:
                    col_names = [col.get('name', '') for col in columns[:5]]
                    output_info.append(f"  - 주요 컬럼: {', '.join(col_names)}{'...' if len(columns) > 5 else ''}")
            elif output_data.get('type') == 'StatisticsOutput':
                output_info.append(f"• {port_name} ({port_type}): 통계 분석 결과")
            elif output_data.get('type') == 'TrainedModelOutput':
                output_info.append(f"• {port_name} ({port_type}): 훈련된 모델")
            elif output_data.get('type') == 'StatsModelsResultOutput':
                output_info.append(f"• {port_name} ({port_type}): 통계 모델 결과")
            elif output_data.get('type') == 'EvaluationOutput':
                output_info.append(f"• {port_name} ({port_type}): 모델 평가 결과")
            else:
                output_info.append(f"• {port_name} ({port_type}): {output_data.get('type', '알 수 없는 타입')}")
        else:
            output_info.append(f"• {port_name} ({port_type}): 실행되지 않음")
    
    return '\n'.join(output_info) if output_info else "출력 데이터 없음"

def get_model_equation(output_data):
    """ResultModel의 함수식 생성"""
    if not output_data or output_data.get('type') != 'StatsModelsResultOutput':
        return None
    
    model_type = output_data.get('modelType', 'Unknown')
    summary = output_data.get('summary', {})
    coefficients = summary.get('coefficients', {})
    feature_columns = output_data.get('featureColumns', [])
    label_column = output_data.get('labelColumn', 'y')
    
    if not coefficients:
        return None
    
    # 절편(intercept) 찾기
    intercept = coefficients.get('const', {}).get('coef', 0)
    
    # 함수식 생성
    equation_parts = []
    for feature in feature_columns:
        coef_info = coefficients.get(feature, {})
        coef = coef_info.get('coef', 0)
        if coef != 0:
            if coef > 0:
                equation_parts.append(f"+ {coef:.4f}×{feature}")
            else:
                equation_parts.append(f"{coef:.4f}×{feature}")
    
    equation = f"{label_column} = {intercept:.4f}"
    if equation_parts:
        equation += " " + " ".join(equation_parts)
    
    return equation

def get_analysis_description(module, all_modules=None, connections=None):
    """모듈의 분석 내용 설명 생성"""
    module_type = module.get('type', 'Unknown')
    module_name = module.get('name', module_type)
    parameters = module.get('parameters', {})
    output_data = module.get('outputData')
    
    analysis_parts = []
    
    # 모듈 타입에 따른 분석 내용
    if module_type == 'LoadData':
        source = parameters.get('source', '알 수 없음')
        analysis_parts.append(f"데이터 소스: {source}")
        analysis_parts.append("CSV 파일에서 데이터를 로드하여 데이터프레임으로 변환")
        analysis_parts.append("데이터 흐름: 파일 → 데이터프레임 변환 → 다음 모듈로 전달")
    elif module_type == 'SelectData':
        selected_cols = parameters.get('columnSelections', {})
        if selected_cols:
            selected_list = [k for k, v in selected_cols.items() if v]
            analysis_parts.append(f"선택된 컬럼: {len(selected_list)}개")
            if selected_list:
                col_names = ', '.join(selected_list[:10])
                if len(selected_list) > 10:
                    col_names += f" 외 {len(selected_list) - 10}개"
                analysis_parts.append(f"  - {col_names}")
        else:
            # 모든 컬럼 선택
            if output_data and output_data.get('type') == 'DataPreview':
                columns = output_data.get('columns', [])
                analysis_parts.append(f"전체 컬럼 선택: {len(columns)}개")
        analysis_parts.append("분석 방법: 입력 데이터에서 필요한 컬럼만 필터링하여 출력")
        analysis_parts.append("데이터 흐름: 전체 데이터 → 컬럼 선택 → 선택된 데이터 출력")
    elif module_type == 'HandleMissingValues':
        method = parameters.get('method', 'unknown')
        if method == 'remove_row':
            analysis_parts.append("결측값 처리 방법: 행 제거")
            analysis_parts.append("분석 방법: 결측값이 포함된 행을 완전히 제거하여 완전한 데이터만 유지")
        elif method == 'impute':
            strategy = parameters.get('strategy', 'mean')
            analysis_parts.append(f"결측값 처리 방법: 대체 ({strategy})")
            analysis_parts.append(f"분석 방법: 결측값을 {strategy} 값으로 대체하여 데이터 손실 최소화")
        elif method == 'knn':
            n_neighbors = parameters.get('n_neighbors', 5)
            analysis_parts.append(f"결측값 처리 방법: KNN 기반 대체 (n_neighbors={n_neighbors})")
            analysis_parts.append(f"분석 방법: 가장 가까운 {n_neighbors}개 이웃의 값을 사용하여 결측값 예측")
        analysis_parts.append("데이터 흐름: 입력 데이터 → 결측값 검출 → 처리 적용 → 정제된 데이터 출력")
    elif module_type == 'Statistics':
        analysis_parts.append("분석 방법: 기술 통계량 계산")
        analysis_parts.append("- 평균, 표준편차, 최소값, 최대값, 사분위수 등")
        analysis_parts.append("- 각 컬럼별 분포 및 요약 통계 제공")
        analysis_parts.append("데이터 흐름: 입력 데이터 → 통계량 계산 → 요약 결과 출력")
    elif module_type in ['LinearRegression', 'LogisticRegression', 'PoissonRegression']:
        analysis_parts.append(f"{module_type} 모델 정의")
        analysis_parts.append("머신러닝 모델 구조 설정")
    elif module_type == 'TrainModel':
        analysis_parts.append("모델 훈련 수행")
        analysis_parts.append("학습 데이터로 모델 파라미터 최적화")
    elif module_type == 'ResultModel':
        feature_cols = parameters.get('feature_columns', [])
        label_col = parameters.get('label_column', '')
        
        # 연결된 모델 정의 모듈 찾기
        model_type = 'Unknown'
        if all_modules and connections:
            for conn in connections:
                if conn.get('to', {}).get('moduleId') == module.get('id') and \
                   conn.get('to', {}).get('portName') == 'model_in':
                    from_module_id = conn.get('from', {}).get('moduleId')
                    model_module = next((m for m in all_modules if m.get('id') == from_module_id), None)
                    if model_module:
                        if model_module.get('type') == 'StatModels':
                            model_type = model_module.get('parameters', {}).get('model', 'Unknown')
                        else:
                            model_type = model_module.get('type', 'Unknown')
                        break
        
        analysis_parts.append(f"사용 모델: {model_type}")
        analysis_parts.append(f"특성 변수: {len(feature_cols)}개 - {', '.join(feature_cols[:5])}{'...' if len(feature_cols) > 5 else ''}")
        analysis_parts.append(f"목표 변수: {label_col}")
        
        # outputData에서 모델 결과 정보 추출
        if output_data and output_data.get('type') == 'StatsModelsResultOutput':
            actual_model_type = output_data.get('modelType', model_type)
            summary = output_data.get('summary', {})
            metrics = summary.get('metrics', {})
            
            analysis_parts.append(f"적용된 모델: {actual_model_type}")
            analysis_parts.append("분석 방법:")
            if actual_model_type == 'OLS':
                analysis_parts.append("  - 최소제곱법(OLS)을 사용한 선형 회귀")
            elif actual_model_type in ['Logistic', 'Logit']:
                analysis_parts.append("  - 로지스틱 회귀를 사용한 이항 분류")
            elif actual_model_type == 'Poisson':
                analysis_parts.append("  - 포아송 회귀를 사용한 카운트 데이터 모델링")
            elif actual_model_type == 'QuasiPoisson':
                analysis_parts.append("  - 준포아송 회귀를 사용한 과분산 카운트 데이터 모델링")
            elif actual_model_type == 'NegativeBinomial':
                analysis_parts.append("  - 음이항 회귀를 사용한 과분산 카운트 데이터 모델링")
            
            # 함수식 생성
            equation = get_model_equation(output_data)
            if equation:
                analysis_parts.append(f"생성된 함수식: {equation}")
            
            # 주요 지표
            if metrics:
                key_metrics = []
                if 'R-squared' in metrics:
                    key_metrics.append(f"R² = {metrics['R-squared']}")
                if 'AIC' in metrics:
                    key_metrics.append(f"AIC = {metrics['AIC']}")
                if 'Log-Likelihood' in metrics:
                    key_metrics.append(f"Log-Likelihood = {metrics['Log-Likelihood']}")
                if key_metrics:
                    analysis_parts.append(f"주요 지표: {', '.join(key_metrics)}")
        else:
            analysis_parts.append("분석 방법: 통계 모델 피팅 및 결과 분석")
        
        analysis_parts.append("데이터 흐름: 데이터 + 모델 정의 → 모델 피팅 → 계수 및 통계량 계산 → 결과 출력")
    elif module_type == 'PredictModel':
        analysis_parts.append("분석 방법: 훈련된 모델을 사용하여 예측 수행")
        analysis_parts.append("예측 과정: 입력 데이터에 모델 함수식 적용 → 예측값 계산 → 결과 데이터에 예측값 컬럼 추가")
        analysis_parts.append("데이터 흐름: 데이터 + 훈련된 모델 → 예측 수행 → 예측값 포함 데이터 출력")
    elif module_type == 'StatModels':
        model_type = parameters.get('model', 'Unknown')
        analysis_parts.append(f"모델 정의: {model_type} 모델 구조 설정")
        analysis_parts.append("분석 방법: 통계 모델의 구조와 파라미터를 정의하여 모델 인스턴스 생성")
        analysis_parts.append("데이터 흐름: 모델 타입 선택 → 모델 정의 생성 → Result Model로 전달")
    elif module_type == 'TrainModel':
        analysis_parts.append("분석 방법: 모델 훈련 수행")
        analysis_parts.append("훈련 과정: 학습 데이터로 모델 파라미터 최적화 → 손실 함수 최소화 → 최적 파라미터 도출")
        analysis_parts.append("데이터 흐름: 데이터 + 모델 정의 → 파라미터 최적화 → 훈련된 모델 출력")
    else:
        analysis_parts.append(f"{module_type} 모듈 실행")
        if parameters:
            key_params = list(parameters.keys())[:3]
            analysis_parts.append(f"주요 파라미터: {', '.join(key_params)}")
        analysis_parts.append("데이터 흐름: 입력 데이터 처리 → 분석 수행 → 결과 출력")
    
    return '\n'.join(analysis_parts)

def create_flowchart_slide(prs, modules, connections, project_name):
    """전체 모델 흐름도를 그리는 슬라이드 생성"""
    # 빈 슬라이드 생성
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃
    
    # 제목 추가
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_shape.text_frame
    title_frame.text = f"전체 모델 흐름도: {project_name}"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(20)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 0, 0)  # 검정색
    title_para.alignment = PP_ALIGN.CENTER
    
    # 유효한 모듈만 필터링
    valid_modules = [m for m in modules if m.get('type') not in EXCLUDED_TYPES]
    
    if not valid_modules:
        return
    
    # 모듈 위치 계산 (간단한 레이아웃: 왼쪽에서 오른쪽으로)
    # 각 모듈을 노드로 표현 (크기 증가)
    module_positions = {}
    module_width = Inches(2.0)  # 1.5 -> 2.0으로 증가
    module_height = Inches(1.0)  # 0.8 -> 1.0으로 증가
    start_x = Inches(0.3)  # 좌측 여백 줄임
    start_y = Inches(1.2)  # 상단 여백 줄임
    spacing_x = Inches(2.5)  # 2.2 -> 2.5로 증가 (더 넓게 배치)
    spacing_y = Inches(1.5)  # 1.2 -> 1.5로 증가
    
    # 간단한 레이어링: 입력이 없는 모듈부터 배치
    def get_module_level(module_id, visited=None):
        if visited is None:
            visited = set()
        if module_id in visited:
            return 0
        visited.add(module_id)
        
        # 이 모듈로 들어오는 연결 찾기
        incoming = [c for c in connections if c.get('to', {}).get('moduleId') == module_id]
        if not incoming:
            return 0
        
        # 최대 깊이 계산
        max_depth = 0
        for conn in incoming:
            from_id = conn.get('from', {}).get('moduleId')
            depth = get_module_level(from_id, visited.copy())
            max_depth = max(max_depth, depth)
        
        return max_depth + 1
    
    # 모듈들을 레벨별로 그룹화
    modules_by_level = {}
    for module in valid_modules:
        level = get_module_level(module.get('id'))
        if level not in modules_by_level:
            modules_by_level[level] = []
        modules_by_level[level].append(module)
    
    # 각 레벨의 모듈들을 배치
    max_modules_per_level = max(len(modules_by_level.get(level, [])) for level in modules_by_level.keys()) if modules_by_level else 1
    
    for level in sorted(modules_by_level.keys()):
        level_modules = modules_by_level[level]
        x = start_x + level * spacing_x
        y_start = start_y
        
        # 레벨 내에서 모듈들을 세로로 배치
        for idx, module in enumerate(level_modules):
            y = y_start + idx * spacing_y
            module_positions[module.get('id')] = {
                'x': x,
                'y': y,
                'module': module
            }
    
    # 모듈 박스 그리기
    for module_id, pos_info in module_positions.items():
        module = pos_info['module']
        x = pos_info['x']
        y = pos_info['y']
        
        # 박스 그리기
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, module_width, module_height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(230, 240, 255)
        box.line.color.rgb = RGBColor(100, 150, 200)
        box.line.width = Pt(1.5)
        
        # 모듈 이름 텍스트
        text_frame = box.text_frame
        text_frame.text = module.get('name', module.get('type'))
        text_frame.word_wrap = True
        para = text_frame.paragraphs[0]
        para.font.size = Pt(12)  # 9 -> 12로 증가
        para.font.bold = True
        para.font.color.rgb = RGBColor(0, 0, 0)  # 검정색
        para.alignment = PP_ALIGN.CENTER
        
        # 모듈 타입을 작은 텍스트로 추가
        if len(text_frame.paragraphs) == 1:
            p = text_frame.add_paragraph()
            p.text = f"({module.get('type')})"
            p.font.size = Pt(10)  # 7 -> 10으로 증가
            p.font.color.rgb = RGBColor(0, 0, 0)  # 검정색
            p.alignment = PP_ALIGN.CENTER
    
    # 연결선 그리기 (add_line 사용 - 더 안정적)
    for conn in connections:
        from_id = conn.get('from', {}).get('moduleId')
        to_id = conn.get('to', {}).get('moduleId')
        
        if from_id not in module_positions or to_id not in module_positions:
            continue
        
        from_pos = module_positions[from_id]
        to_pos = module_positions[to_id]
        
        # 시작점과 끝점 계산 (박스의 오른쪽 중앙과 왼쪽 중앙)
        start_x = from_pos['x'] + module_width
        start_y = from_pos['y'] + module_height / 2
        end_x = to_pos['x']
        end_y = to_pos['y'] + module_height / 2
        
        # 선 그리기 (add_line 사용 - 더 안정적이고 호환성 좋음)
        try:
            # 일반 선 사용 (connector 대신)
            line = slide.shapes.add_line(start_x, start_y, end_x, end_y)
            line.line.color.rgb = RGBColor(100, 100, 100)
            line.line.width = Pt(2)
            
            # 화살표 끝 스타일 설정
            from pptx.enum.dml import MSO_ARROWHEAD_LENGTH, MSO_ARROWHEAD_WIDTH, MSO_ARROWHEAD_TYPE
            line.line.end_arrowhead_length = MSO_ARROWHEAD_LENGTH.MEDIUM
            line.line.end_arrowhead_width = MSO_ARROWHEAD_WIDTH.MEDIUM
            line.line.end_arrowhead_type = MSO_ARROWHEAD_TYPE.TRIANGLE
        except Exception as e:
            # 화살표 설정이 실패해도 기본 선은 그리기
            try:
                line = slide.shapes.add_line(start_x, start_y, end_x, end_y)
                line.line.color.rgb = RGBColor(100, 100, 100)
                line.line.width = Pt(2)
            except:
                # 최후의 수단: 에러 무시하고 계속 진행
                print(f"연결선 그리기 실패: {e}")
                pass

def create_module_slide(prs, module, all_modules, connections, module_index, total_modules):
    """하나의 모듈에 대한 슬라이드 생성"""
    module_type = module.get('type', 'Unknown')
    module_name = module.get('name', module_type)
    
    # 슬라이드 추가 (빈 레이아웃 사용)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃
    
    # 제목을 좌측 상단에 작은 공간으로 배치
    title_shape = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(2.5), Inches(0.5))
    title_frame = title_shape.text_frame
    title_frame.text = f"{module_index}. {module_name}"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(18)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.LEFT
    
    # 본문 영역을 넓혀서 작성 (제목 아래부터 시작, 좌우 마진을 줄임)
    content_shape = slide.shapes.add_textbox(Inches(0.3), Inches(0.8), Inches(9.4), Inches(6.5))
    tf = content_shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    
    # 모든 텍스트의 기본 폰트 크기 설정
    def set_font_size(paragraph, size):
        paragraph.font.size = Pt(size)
        for run in paragraph.runs:
            run.font.size = Pt(size)
    
    # 입력 데이터 형태
    p = tf.paragraphs[0]
    p.text = "📥 입력 데이터 형태"
    p.level = 0
    p.font.size = Pt(14)  # 13 -> 14로 변경
    p.font.bold = True
    p.space_after = Pt(6)
    
    input_info = get_input_data_info(module, all_modules, connections)
    p = tf.add_paragraph()
    p.text = input_info
    p.level = 1
    p.font.size = Pt(13)  # 11 -> 13으로 변경
    p.space_after = Pt(8)
    
    # 분석 내용
    p = tf.add_paragraph()
    p.text = "🔍 분석 내용"
    p.level = 0
    p.font.size = Pt(14)  # 13 -> 14로 변경
    p.font.bold = True
    p.space_before = Pt(6)
    p.space_after = Pt(6)
    
    analysis_desc = get_analysis_description(module, all_modules, connections)
    # 분석 내용을 여러 줄로 분리
    for line in analysis_desc.split('\n'):
        p = tf.add_paragraph()
        p.text = line
        p.level = 1
        p.font.size = Pt(13)  # 11 -> 13으로 변경
        p.space_after = Pt(3)
    
    # 생성되는 결과 파일
    p = tf.add_paragraph()
    p.text = "📤 생성되는 결과"
    p.level = 0
    p.font.size = Pt(14)  # 13 -> 14로 변경
    p.font.bold = True
    p.space_before = Pt(6)
    p.space_after = Pt(6)
    
    output_info = get_output_data_info(module)
    for line in output_info.split('\n'):
        p = tf.add_paragraph()
        p.text = line
        p.level = 1
        p.font.size = Pt(13)  # 11 -> 13으로 변경
        p.space_after = Pt(3)
    
    # 파라미터 정보 (있는 경우, 공간이 남으면)
    parameters = module.get('parameters', {})
    if parameters and len(parameters) > 0 and len(tf.paragraphs) < 20:  # 공간이 충분할 때만
        p = tf.add_paragraph()
        p.text = "⚙️ 주요 파라미터"
        p.level = 0
        p.font.size = Pt(14)  # 13 -> 14로 변경
        p.font.bold = True
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        
        for key, value in list(parameters.items())[:3]:  # 최대 3개만 표시
            p = tf.add_paragraph()
            if isinstance(value, (dict, list)):
                value_str = json.dumps(value, ensure_ascii=False)[:40] + "..."
            else:
                value_str = str(value)[:40]
            p.text = f"  • {key}: {value_str}"
            p.level = 1
            p.font.size = Pt(13)  # 11 -> 13으로 변경
            p.space_after = Pt(2)

def create_single_ppt(project_data, output_path):
    """하나의 PPT 파일에 모든 모듈 슬라이드 생성"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    modules = project_data.get('modules', [])
    connections = project_data.get('connections', [])
    project_name = project_data.get('projectName', 'Untitled Project')
    
    # TextBox, GroupBox 같은 Shape 타입은 제외
    valid_modules = [m for m in modules if m.get('type') not in EXCLUDED_TYPES]
    
    if not valid_modules:
        print("유효한 모듈이 없습니다.")
        return None
    
    # 제목 슬라이드
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    title.text = project_name
    subtitle.text = f"총 {len(valid_modules)}개 모듈"
    
    # 전체 흐름도 슬라이드 추가
    print("전체 흐름도 슬라이드 생성 중...")
    create_flowchart_slide(prs, modules, connections, project_name)
    
    # 각 모듈에 대한 슬라이드 생성
    print(f"총 {len(valid_modules)}개의 모듈에 대해 슬라이드를 생성합니다...")
    for i, module in enumerate(valid_modules, 1):
        try:
            module_name = module.get('name', module.get('type'))
            print(f"[{i}/{len(valid_modules)}] {module_name} 슬라이드 생성 중...")
            create_module_slide(prs, module, modules, connections, i, len(valid_modules))
        except Exception as e:
            print(f"모듈 {module.get('name', module.get('type'))} 슬라이드 생성 실패: {e}")
            import traceback
            traceback.print_exc()
    
    # 파일 저장
    prs.save(output_path)
    print(f"\nPPT 파일 생성 완료: {output_path}")
    print(f"다운로드 폴더에 저장되었습니다: {os.path.dirname(output_path)}")
    
    return output_path

def process_project_data(project_data, output_dir=None):
    """프로젝트 데이터를 읽어서 하나의 PPT 파일 생성"""
    # 출력 디렉토리 설정 (다운로드 폴더 또는 지정된 폴더)
    if output_dir is None:
        output_dir = get_download_folder()
        print(f"다운로드 폴더에 저장: {output_dir}")
    else:
        os.makedirs(output_dir, exist_ok=True)
        print(f"지정된 폴더에 저장: {output_dir}")
    
    # 다운로드 폴더가 존재하는지 확인
    if not os.path.exists(output_dir):
        print(f"경고: 다운로드 폴더가 존재하지 않습니다. 생성합니다: {output_dir}")
        os.makedirs(output_dir, exist_ok=True)
    
    project_name = project_data.get('projectName', 'Untitled_Project')
    safe_name = sanitize_filename(project_name)
    filename = f"{safe_name}_모듈분석.pptx"
    output_path = os.path.join(output_dir, filename)
    print(f"저장 경로: {output_path}")
    
    # PPT 파일 생성
    result_path = create_single_ppt(project_data, output_path)
    
    if result_path:
        return [{
            'filename': filename,
            'filepath': result_path,
            'module_count': len([m for m in project_data.get('modules', []) if m.get('type') not in EXCLUDED_TYPES])
        }]
    else:
        return []

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("사용법: python generate_module_ppts.py <project_json_file> [output_dir]")
        print("  output_dir을 지정하지 않으면 다운로드 폴더에 저장됩니다.")
        sys.exit(1)
    
    input_file = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else None
    
    if not os.path.exists(input_file):
        print(f"파일을 찾을 수 없습니다: {input_file}")
        sys.exit(1)
    
    try:
        with open(input_file, 'r', encoding='utf-8') as f:
            project_data = json.load(f)
        
        generated_files = process_project_data(project_data, output_dir)
        
        if generated_files:
            print(f"\n생성 완료: {generated_files[0]['filepath']}")
            print(f"모듈 수: {generated_files[0]['module_count']}개")
        else:
            print("\nPPT 파일 생성에 실패했습니다.")
        
    except Exception as e:
        print(f"오류 발생: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
